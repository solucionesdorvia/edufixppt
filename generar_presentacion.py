#!/usr/bin/env python3
"""
Genera presentación comercial EduFix para proyector (pitch de venta).
Lee solo respuestas.csv — no modifica el CSV.
Salida: EduFix_Pitch_Venta.pptx
"""
from __future__ import annotations

import csv
from collections import Counter
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).resolve().parent
CSV_PATH = BASE / "respuestas.csv"
OUT_PATH = BASE / "EduFix_Pitch_Venta.pptx"

C_NAVY = RGBColor(0x0A, 0x1C, 0x36)
C_ORANGE = RGBColor(0xD9, 0x53, 0x1E)
C_TEAL = RGBColor(0x1F, 0x8B, 0x63)
C_BLUE = RGBColor(0x1F, 0x57, 0x8B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED = RGBColor(0xC8, 0xD4, 0xE8)
C_CARD = RGBColor(0x12, 0x24, 0x42)
C_ACCENT_LIGHT = RGBColor(0x4E, 0xD4, 0x9E)


def load_rows(path: Path) -> list[dict]:
    with open(path, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def mean(nums: list[float]) -> float:
    return sum(nums) / len(nums) if nums else 0.0


def pct(part: int, total: int) -> float:
    return 100.0 * part / total if total else 0.0


def fmt_num(x: float, dec: int = 1) -> str:
    return f"{x:.{dec}f}".replace(".", ",")


def add_deck_chrome(slide, prs: Presentation) -> None:
    h, w = prs.slide_height, prs.slide_width
    stripe = Emu(int(w * 0.014))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, stripe, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ORANGE
    bar.line.fill.background()
    fh = Emu(int(h * 0.014))
    third = w // 3
    for i, col in enumerate((C_ORANGE, C_TEAL, C_BLUE)):
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, third * i, h - fh, third, fh)
        seg.fill.solid()
        seg.fill.fore_color.rgb = col
        seg.line.fill.background()


def tb(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size_pt: float = 14,
    bold: bool = False,
    color: RGBColor = C_WHITE,
    align=PP_ALIGN.LEFT,
    font_name: str = "Calibri",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name


def slide_dark(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_NAVY
    return slide


def add_finding_cards(slide, prs, findings: list[tuple[str, str, RGBColor]]) -> None:
    w, h = prs.slide_width, prs.slide_height
    pad = Emu(int(w * 0.05))
    top0 = Emu(int(h * 0.2))
    gap_x = Emu(int(w * 0.025))
    gap_y = Emu(int(h * 0.032))
    cell_w = (w - 2 * pad - gap_x) // 2
    cell_h = Emu(int(h * 0.195))

    for i, (stat, desc, accent) in enumerate(findings):
        row, col = divmod(i, 2)
        left = pad + col * (cell_w + gap_x)
        top = top0 + row * (cell_h + gap_y)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cell_w, cell_h)
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD
        card.line.color.rgb = RGBColor(0x35, 0x50, 0x70)
        card.line.width = Pt(0.75)

        ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(int(w * 0.011)), cell_h)
        ab.fill.solid()
        ab.fill.fore_color.rgb = accent
        ab.line.fill.background()

        tb(
            slide,
            left + Emu(int(w * 0.038)),
            top + Emu(int(h * 0.014)),
            cell_w - Emu(int(w * 0.048)),
            Emu(int(cell_h * 0.44)),
            stat,
            size_pt=34,
            bold=True,
            color=C_WHITE,
        )
        tb(
            slide,
            left + Emu(int(w * 0.038)),
            top + Emu(int(cell_h * 0.48)),
            cell_w - Emu(int(w * 0.048)),
            Emu(int(cell_h * 0.48)),
            desc,
            size_pt=14,
            bold=False,
            color=C_MUTED,
        )


def pillar_row(
    slide,
    prs,
    items: list[tuple[str, str]],
    top_inches: float,
) -> None:
    """Tres columnas: título corto + texto (para proyector)."""
    w = prs.slide_width
    n = len(items)
    gap = Inches(0.35)
    col_w = (w - Inches(1.7) - gap * (n - 1)) // n
    left0 = Inches(0.85)
    for i, (head, sub) in enumerate(items):
        left = left0 + i * (col_w + gap)
        tb(slide, left, Inches(top_inches), col_w, Inches(0.55), head, size_pt=22, bold=True, color=C_ORANGE)
        tb(slide, left, Inches(top_inches + 0.62), col_w, Inches(1.8), sub, size_pt=17, color=C_MUTED)


def chart_series_color(chart, rgb: RGBColor) -> None:
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = rgb


def build_presentation(rows: list[dict]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    n = len(rows)
    if n == 0:
        raise ValueError("No hay filas en respuestas.csv")

    def pc(fn) -> int:
        return sum(1 for r in rows if fn(r))

    pct_no_supo = pct(pc(lambda r: r.get("detectaste_problema_y_no_supiste_a_quien_reportar") == "Si"), n)
    pct_ignora = pct(pc(lambda r: r.get("ante_desperfecto_que_sueles_hacer") == "Ignorarlo"), n)
    pct_no_canal = pct(pc(lambda r: r.get("tiempo_reportar_canales_oficiales") == "No se como hacerlo"), n)
    avg_impacto = mean([float(r["afecta_entorno_descuidado_1_10"]) for r in rows])
    pct_1a3 = pct(pc(lambda r: r.get("veces_por_mes_notas_desperfecto") == "1 a 3 veces"), n)
    pct_edufix = pct(pc(lambda r: r.get("dispuesto_foto_si_notificacion_al_cerrar") == "Si"), n)
    avg_util = mean([float(r["utilidad_app_foto_estado_1_10"]) for r in rows])

    roles = Counter(r.get("rol_principal", "") for r in rows)
    frust = Counter(r.get("frustracion_reportes_actual", "") for r in rows)
    hist_u = Counter()
    hist_a = Counter()
    for r in rows:
        try:
            hist_u[int(float(r["utilidad_app_foto_estado_1_10"]))] += 1
            hist_a[int(float(r["afecta_entorno_descuidado_1_10"]))] += 1
        except (ValueError, TypeError):
            pass

    today = date.today().strftime("%d/%m/%Y")
    lead_valid = (
        "Más de 400 encuestas confirman un proceso ineficiente e informal."
        if n >= 400
        else f"{n} respuestas confirman un patrón claro de fricción en el reporte."
    )

    # --- 1 Portada pitch ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(1.05), Inches(12), Inches(1.35), "EduFix", size_pt=72, bold=True, color=C_ORANGE)
    tb(
        s,
        Inches(0.85),
        Inches(2.35),
        Inches(11.5),
        Inches(1.1),
        "Reportes de mantenimiento en campus:\nclaros, rápidos y accionables.",
        size_pt=32,
        bold=True,
        color=C_WHITE,
    )
    tb(s, Inches(0.85), Inches(3.85), Inches(11), Inches(0.65), f"Pitch comercial · {n} respuestas validadas · {today}", size_pt=20, color=C_MUTED)

    # --- 2 Gancho (una cifra enorme) ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.75), Inches(11), Inches(0.7), "Primero, el dolor", size_pt=26, bold=True, color=C_TEAL)
    tb(
        s,
        Inches(0.85),
        Inches(1.85),
        Inches(12),
        Inches(2.2),
        f"{fmt_num(pct_ignora)}%",
        size_pt=120,
        bold=True,
        color=C_ORANGE,
        align=PP_ALIGN.CENTER,
    )
    tb(
        s,
        Inches(1.2),
        Inches(4.15),
        Inches(11),
        Inches(1.2),
        "de las personas encuestadas suele ignorar el desperfecto.\n"
        "Sin dato no hay prioridad: el campus degrada y nadie actúa a tiempo.",
        size_pt=24,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    # --- 3 Problema (venta) ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.65), Inches(2.5), Inches(0.55), "01", size_pt=40, bold=True, color=C_ORANGE)
    tb(s, Inches(0.85), Inches(1.2), Inches(11), Inches(0.85), "El problema no es la rotura: es el sistema de reporte", size_pt=36, bold=True)
    body = (
        f"{lead_valid}\n\n"
        "Hoy el flujo es informal: correos dispersos, pasillos, incertidumbre sobre el estado del arreglo.\n"
        "Eso genera desconfianza, retrabajo y una mala experiencia en el espacio donde estudian y enseñan."
    )
    tb(s, Inches(0.85), Inches(2.25), Inches(11.3), Inches(4.5), body, size_pt=22, color=C_MUTED)

    # --- 4 La evidencia ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(1.0), Inches(11), Inches(1.0), "Lo medimos con datos", size_pt=40, bold=True)
    tb(
        s,
        Inches(0.85),
        Inches(2.2),
        Inches(11.2),
        Inches(2.5),
        f"{n} respuestas reales sobre percepción, frustración y disposición a cambiar de hábito.\n\n"
        "Nada de suposiciones: números para presentar a dirección, mantenimiento y comunidad educativa.",
        size_pt=24,
        color=C_MUTED,
    )

    # --- 5 Hallazgos 6 ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.45), Inches(2.5), Inches(0.45), "03", size_pt=36, bold=True, color=C_ORANGE)
    tb(s, Inches(0.85), Inches(0.95), Inches(11), Inches(0.65), "Hallazgos que abren la conversación", size_pt=32, bold=True)
    tb(s, Inches(0.85), Inches(1.5), Inches(11), Inches(0.4), "Indicadores calculados desde tu dataset (solo lectura).", size_pt=15, color=C_MUTED)
    findings = [
        (f"{fmt_num(pct_no_supo)}%", "No supo a quién reportar", C_ORANGE),
        (f"{fmt_num(pct_ignora)}%", "Ignoró el incidente", RGBColor(0xE8, 0x5D, 0x5D)),
        (f"{fmt_num(pct_no_canal)}%", "No sabe usar el canal oficial", RGBColor(0x5A, 0xB0, 0xFF)),
        (f"{fmt_num(avg_impacto)}/10", "Impacto en la experiencia", C_TEAL),
        (f"{fmt_num(pct_1a3)}%", "Ve 1–3 desperfectos / mes", C_BLUE),
        (f"{fmt_num(pct_edufix)}%", "Usaría EduFix en ~30 s", C_ACCENT_LIGHT),
    ]
    add_finding_cards(s, prs, findings)

    # --- 6 Hoy vs mañana ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.6), Inches(11), Inches(0.75), "Coste de no actuar vs oportunidad", size_pt=34, bold=True)
    mid = Inches(6.55)
    bx_w = Inches(5.9)
    tb(s, Inches(0.85), Inches(1.45), bx_w, Inches(0.45), "Hoy", size_pt=24, bold=True, color=C_ORANGE)
    tb(
        s,
        Inches(0.85),
        Inches(1.95),
        bx_w,
        Inches(4.5),
        "• Reporte opaco y lento\n• Datos no estandarizados\n• Poca trazabilidad para el usuario\n• Mantenimiento reacciona tarde",
        size_pt=20,
        color=C_MUTED,
    )
    tb(s, mid, Inches(1.45), bx_w, Inches(0.45), "Con EduFix", size_pt=24, bold=True, color=C_TEAL)
    tb(
        s,
        mid,
        Inches(1.95),
        bx_w,
        Inches(4.5),
        "• 30 segundos para reportar\n• Foto + contexto estructurado\n• Estado visible para quien reporta\n• Cola priorizada para operaciones",
        size_pt=20,
        color=C_MUTED,
    )

    # --- 7 Solución ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(1.0), Inches(2), Inches(0.7), "04", size_pt=44, bold=True, color=C_ORANGE)
    tb(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(1.1), "Solución", size_pt=52, bold=True)
    tb(
        s,
        Inches(0.85),
        Inches(3.15),
        Inches(11.2),
        Inches(2.2),
        "EduFix es el canal único y móvil para reportar desperfectos:\n"
        "menos fricción, más señal para mantenimiento y mejor experiencia de campus.",
        size_pt=26,
        color=C_MUTED,
    )

    # --- 8 Pilares ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.55), Inches(11), Inches(0.7), "Por qué compran EduFix", size_pt=34, bold=True)
    pillar_row(
        s,
        prs,
        [
            ("Rapidez", "Flujo guiado. El usuario no pierde tiempo en formularios interminables."),
            ("Claridad", "Cada reporte llega con contexto homogéneo: menos idas y vueltas."),
            ("Trazabilidad", "Quien reporta ve que el caso existe y avanza: confianza institucional."),
        ],
        1.35,
    )

    # --- 9 Cómo funciona ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.55), Inches(11), Inches(0.7), "Cómo funciona (30 segundos)", size_pt=34, bold=True)
    pillar_row(
        s,
        prs,
        [
            ("1 · Captura", "Foto y categoría del problema en el lugar del hecho."),
            ("2 · Envío", "Un toque: el reporte entra a la cola oficial con geolocalización / área."),
            ("3 · Cierre de ciclo", "Notificación cuando hay avance o resolución."),
        ],
        1.35,
    )

    # --- 10 Stakeholders ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.55), Inches(11), Inches(0.7), "Valor para cada actor", size_pt=34, bold=True)
    pillar_row(
        s,
        prs,
        [
            ("Alumnos y docentes", "Menos fricción; sensación de campus cuidado."),
            ("Mantenimiento", "Tickets ordenados, con evidencia y menos ruido."),
            ("Dirección", "Métricas de demanda real y cumplimiento visible."),
        ],
        1.35,
    )

    # --- 11 Roles ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.5), Inches(11), Inches(0.75), "¿Quién respondió? (muestra)", size_pt=30, bold=True)
    chart_data = CategoryChartData()
    cats = [k for k, _ in roles.most_common()]
    chart_data.categories = cats
    chart_data.add_series("n", tuple(roles[c] for c in cats))
    ch = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(1.35), Inches(11.4), Inches(5.35), chart_data
    ).chart
    ch.has_legend = False
    ch.plots[0].has_data_labels = True
    chart_series_color(ch, C_ORANGE)

    # --- 12 Frustración ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.5), Inches(11), Inches(0.75), "Qué duele del proceso actual", size_pt=30, bold=True)
    chart_data = CategoryChartData()
    cats = [k for k, _ in frust.most_common()]
    chart_data.categories = cats
    chart_data.add_series("n", tuple(frust[c] for c in cats))
    ch = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(1.35), Inches(11.4), Inches(5.35), chart_data
    ).chart
    ch.has_legend = False
    ch.plots[0].has_data_labels = True
    chart_series_color(ch, C_TEAL)

    # --- 13 Cierre estadístico (venta) ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.65), Inches(11), Inches(0.65), "El número que cierra la reunión", size_pt=26, bold=True, color=C_TEAL)
    tb(
        s,
        Inches(0.85),
        Inches(1.55),
        Inches(12),
        Inches(2.0),
        f"{fmt_num(pct_edufix)}%",
        size_pt=110,
        bold=True,
        color=C_ACCENT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    tb(
        s,
        Inches(1.0),
        Inches(3.85),
        Inches(11),
        Inches(1.3),
        "declaran que usarían EduFix si el reporte tomara unos 30 segundos.\n"
        "La demanda de una herramienta simple ya está validada.",
        size_pt=24,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    # --- 14 Utilidad + impacto (compacto) ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.5), Inches(11), Inches(0.65), "Percepción de utilidad e impacto (1–10)", size_pt=28, bold=True)
    half_w = Inches(5.55)
    chart_data = CategoryChartData()
    chart_data.categories = [str(i) for i in range(1, 11)]
    chart_data.add_series("u", tuple(hist_u.get(i, 0) for i in range(1, 11)))
    ch1 = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.85), Inches(1.25), half_w, Inches(4.9), chart_data
    ).chart
    ch1.has_legend = False
    ch1.plots[0].has_data_labels = True
    chart_series_color(ch1, C_ORANGE)
    tb(s, Inches(0.85), Inches(1.05), half_w, Inches(0.35), f"Utilidad (media {fmt_num(avg_util)})", size_pt=16, bold=True, color=C_MUTED)

    chart_data = CategoryChartData()
    chart_data.categories = [str(i) for i in range(1, 11)]
    chart_data.add_series("a", tuple(hist_a.get(i, 0) for i in range(1, 11)))
    x2 = Inches(0.85) + half_w + Inches(0.35)
    ch2 = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x2, Inches(1.25), half_w, Inches(4.9), chart_data
    ).chart
    ch2.has_legend = False
    ch2.plots[0].has_data_labels = True
    chart_series_color(ch2, C_TEAL)
    tb(s, x2, Inches(1.05), half_w, Inches(0.35), f"Impacto entorno (media {fmt_num(avg_impacto)})", size_pt=16, bold=True, color=C_MUTED)

    # --- 15 CTA ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(0.7), Inches(11), Inches(0.85), "Próximo paso", size_pt=38, bold=True)
    tb(
        s,
        Inches(0.85),
        Inches(1.65),
        Inches(11.3),
        Inches(4.8),
        "• Piloto en un edificio o facultad (4–6 semanas).\n"
        "• Definir integración con mesa de ayuda / mantenimiento existente.\n"
        "• Métricas de adopción: reportes / semana, tiempo medio de cierre, satisfacción.\n\n"
        "Propuesta: demo en vivo + acceso al panel de datos que alimenta esta presentación.",
        size_pt=22,
        color=C_MUTED,
    )

    # --- 16 Cierre ---
    s = slide_dark(prs)
    add_deck_chrome(s, prs)
    tb(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(1.2), "Gracias", size_pt=56, bold=True, align=PP_ALIGN.CENTER)
    tb(
        s,
        Inches(0.85),
        Inches(3.25),
        Inches(11.5),
        Inches(2.0),
        "EduFix\n"
        "Reportes que se ven, se priorizan y se resuelven.\n\n"
        f"{n} respuestas · {today}",
        size_pt=22,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    return prs


def main():
    if not CSV_PATH.is_file():
        raise SystemExit(f"No se encontró {CSV_PATH}")
    rows = load_rows(CSV_PATH)
    prs = build_presentation(rows)
    prs.save(OUT_PATH)
    print(f"Guardado: {OUT_PATH} ({len(rows)} filas · {len(prs.slides)} diapositivas · pitch proyector)")


if __name__ == "__main__":
    main()
